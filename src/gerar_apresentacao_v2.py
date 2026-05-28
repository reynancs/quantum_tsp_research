# -*- coding: utf-8 -*-
"""
Apresentacao oral (.pptx) — VERSAO 2 — XI SAPCT / SENAI CIMATEC 2026.

Trabalho: "Aplicacao de Quantum Annealing ao Problema do Caixeiro Viajante:
revisao exploratoria da literatura em logistica".

Diferencas em relacao a versao 1:
  - Estrutura segue os TOPICOS do artigo: Resumo, Introducao, Metodologia,
    Resultados e Discussao, Consideracoes Finais, Agradecimentos e Referencias.
  - Dimensionada para uma APRESENTACAO DE ~10 MINUTOS (12 slides).
  - HIERARQUIA TIPOGRAFICA AMPLIADA — corpo de texto em 14 pt, com titulos,
    subtitulos e rotulos escalonados por boas praticas (escala modular ~1.3).

Mesmo modelo/template do evento e a mesma identidade visual da versao 1.

Execucao:  python src/gerar_apresentacao_v2.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

# ----------------------------------------------------------------------
# Caminhos
# ----------------------------------------------------------------------
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS = os.path.join(ROOT, "docs")
IMG = os.path.join(DOCS, "images")
MEDIA = os.path.join(DOCS, "tmpl_unpacked", "ppt", "media")
TEMPLATE = os.path.join(DOCS, "Template_Apresentacao_Oral_SAPCTSLIDES2026B.pptx")
OUTPUT = os.path.join(DOCS, "Apresentacao_SAPCT_CIMATEC_2026_TSP_Quantum_Annealing_v2.pptx")

LOGO_SAPCT = os.path.join(MEDIA, "image4.png")
LOGO_CIMATEC = os.path.join(MEDIA, "image5.png")

# ----------------------------------------------------------------------
# Paleta de cores (identica a versao 1 / tema do template)
# ----------------------------------------------------------------------
NAVY = RGBColor(0x0E, 0x28, 0x41)
BLUE = RGBColor(0x2E, 0x5C, 0x8A)
TEAL = RGBColor(0x15, 0x60, 0x82)
STEEL = RGBColor(0x5B, 0x7F, 0xB5)
LIGHT = RGBColor(0xEC, 0xF1, 0xF6)
CARD = RGBColor(0xF4, 0xF6, 0xF9)
INK = RGBColor(0x2B, 0x2B, 0x2B)
GRAY = RGBColor(0x5E, 0x66, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xDD, 0xE6, 0xEF)

FONT = "Arial"
FOOTER_TXT = "XI SAPCT  ·  SENAI CIMATEC  ·  Salvador-BA  ·  28 e 29 de Maio de 2026"

# ----------------------------------------------------------------------
# ESCALA TIPOGRAFICA  (boas praticas — corpo 14 pt como base)
#   corpo 14  ->  sub-bullet 13  ->  subtitulo 18  ->  destaque 22
#               ->  titulo 32  |  numeros 30  |  legenda 11  |  rodape 10
# ----------------------------------------------------------------------
SZ_TITLE = 32        # titulo do slide
SZ_KICKER = 14       # rotulo de secao (eyebrow)
SZ_LEAD = 22         # frase de destaque / subtitulo de slide
SZ_SUBHEAD = 18      # cabecalho de cartao / bloco
SZ_BODY = 14         # corpo de texto (base solicitada)
SZ_BODY_SM = 13      # sub-itens e textos de apoio
SZ_CAPTION = 11      # legendas de imagens
SZ_STAT = 30         # numeros de destaque
SZ_STAT_LABEL = 12   # rotulo dos indicadores
SZ_FOOTER = 10       # rodape

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

# ----------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------

def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def clear_placeholders(slide):
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def _style_run(run, size, color, bold=False, italic=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def add_text(slide, left, top, width, height, runs_spec, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=4, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs_spec):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for spec in para:
            text, size, color = spec[0], spec[1], spec[2]
            bold = spec[3] if len(spec) > 3 else False
            italic = spec[4] if len(spec) > 4 else False
            r = p.add_run()
            r.text = text
            _style_run(r, size, color, bold, italic)
    return tb


def add_bullets(slide, left, top, width, height, items, size=SZ_BODY,
                color=INK, gap=8, lead_color=TEAL, line_spacing=1.16):
    """Lista de topicos. items: (texto, nivel, bold?)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        text = it[0]
        level = it[1] if len(it) > 1 else 0
        bold = it[2] if len(it) > 2 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        if level == 0:
            mark = p.add_run()
            mark.text = "▪  "
            _style_run(mark, size, lead_color, True)
            r = p.add_run()
            r.text = text
            _style_run(r, size, color, bold)
        else:
            p.level = 1
            mark = p.add_run()
            mark.text = "      –  "
            _style_run(mark, size - 1, STEEL, False)
            r = p.add_run()
            r.text = text
            _style_run(r, size - 1, GRAY, False)
    return tb


def add_pic_fit(slide, path, box_l, box_t, box_w, box_h, align="center"):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    t = box_t + (box_h - h) // 2
    if align == "center":
        l = box_l + (box_w - w) // 2
    elif align == "left":
        l = box_l
    else:
        l = box_l + (box_w - w)
    return slide.shapes.add_picture(path, l, t, w, h)


def add_rect(slide, left, top, width, height, fill, shape=MSO_SHAPE.RECTANGLE,
             line=None, line_w=1.0, radius=None):
    sp = slide.shapes.add_shape(shape, left, top, width, height)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def add_logos(slide, sapct=True):
    if sapct:
        add_pic_fit(slide, LOGO_SAPCT, Inches(11.02), Inches(0.32),
                    Inches(1.80), Inches(1.04), align="right")
    add_pic_fit(slide, LOGO_CIMATEC, Inches(11.55), Inches(6.86),
                Inches(1.28), Inches(0.52), align="right")


def add_footer(slide, number):
    add_rect(slide, Inches(0.62), Inches(7.02), Inches(7.8), Pt(1.6), TEAL)
    add_text(slide, Inches(0.62), Inches(7.08), Inches(8.0), Inches(0.32),
             [[(FOOTER_TXT, SZ_FOOTER, GRAY)]])
    add_text(slide, Inches(10.55), Inches(7.00), Inches(0.90), Inches(0.34),
             [[("%02d" % number, 11, BLUE, True)]], align=PP_ALIGN.RIGHT)


def add_header(slide, kicker, title):
    """Cabecalho: rotulo de secao + titulo + regua. Hierarquia ampliada."""
    add_rect(slide, Inches(0.62), Inches(0.40), Inches(0.22), Inches(0.32), TEAL)
    add_text(slide, Inches(0.96), Inches(0.38), Inches(8.4), Inches(0.36),
             [[(kicker.upper(), SZ_KICKER, TEAL, True)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.60), Inches(0.74), Inches(10.2), Inches(1.00),
             [[(title, SZ_TITLE, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.0)
    add_rect(slide, Inches(0.64), Inches(1.66), Inches(1.00), Pt(3.6), NAVY)


def content_slide(prs, kicker, title, number, sapct=True):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # BLANK
    clear_placeholders(slide)
    set_bg(slide, WHITE)
    add_header(slide, kicker, title)
    add_logos(slide, sapct=sapct)
    add_footer(slide, number)
    return slide


def caption(slide, left, top, width, text):
    add_text(slide, left, top, width, Inches(0.3),
             [[(text, SZ_CAPTION, GRAY, False, True)]], align=PP_ALIGN.CENTER)


def stat_card(slide, left, top, width, height, big, label, big_color=NAVY,
              fill=LIGHT):
    add_rect(slide, left, top, width, height, fill,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    add_rect(slide, left, top, Inches(0.08), height, TEAL,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    add_text(slide, left + Inches(0.22), top + Inches(0.12),
             width - Inches(0.36), Inches(0.56),
             [[(big, SZ_STAT, big_color, True)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, left + Inches(0.22), top + height - Inches(0.66),
             width - Inches(0.36), Inches(0.58),
             [[(label, SZ_STAT_LABEL, GRAY)]], anchor=MSO_ANCHOR.TOP,
             line_spacing=1.04)


def chip(slide, left, top, width, text):
    h = Inches(0.46)
    add_rect(slide, left, top, width, h, NAVY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    add_text(slide, left, top, width, h, [[(text, SZ_BODY_SM, WHITE, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ======================================================================
# MONTAGEM
# ======================================================================
prs = Presentation(TEMPLATE)

# mantem a capa do evento (slide 1) e remove os 2 slides-modelo
xml_slides = prs.slides._sldIdLst
for sid in list(xml_slides)[1:]:
    rId = sid.get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(sid)

# ----------------------------------------------------------------------
# SLIDE 2 — Titulo do trabalho
# ----------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[0])
clear_placeholders(s)
set_bg(s, WHITE)
add_rect(s, 0, 0, Inches(0.34), EMU_H, NAVY)
add_rect(s, Inches(0.34), 0, Inches(0.06), EMU_H, TEAL)
add_pic_fit(s, LOGO_SAPCT, Inches(0.95), Inches(0.58), Inches(3.0), Inches(1.55),
            align="left")
add_text(s, Inches(0.98), Inches(2.24), Inches(8.0), Inches(0.36),
         [[("REVISÃO EXPLORATÓRIA DA LITERATURA", SZ_KICKER, TEAL, True)]])
add_text(s, Inches(0.98), Inches(2.64), Inches(11.5), Inches(1.9),
         [[("Aplicação de Quantum Annealing ao Problema do "
            "Caixeiro Viajante", 34, NAVY, True)],
          [("Revisão exploratória da literatura em logística",
            SZ_LEAD, BLUE, False, True)]],
         line_spacing=1.05, space_after=8)
add_rect(s, Inches(1.0), Inches(4.40), Inches(1.1), Pt(3.6), TEAL)
add_text(s, Inches(0.98), Inches(4.62), Inches(11.4), Inches(0.5),
         [[("Renan Cardoso dos Santos", 16, INK, True),
           ("¹     ", 12, TEAL, True),
           ("Valéria Loureiro da Silva", 16, INK, True),
           ("²     ", 12, TEAL, True),
           ("Anderson Rafael Correia B. da Silva", 16, INK, True),
           ("³", 12, TEAL, True)]])
add_text(s, Inches(0.98), Inches(5.04), Inches(11.4), Inches(1.0),
         [[("¹ Mestrando em Gestão de Tecnologia e Inovação "
            "— Universidade SENAI CIMATEC, Salvador-BA", 12, GRAY)],
          [("² ³ Universidade SENAI CIMATEC, Salvador-BA", 12, GRAY)]],
         space_after=4)
add_rect(s, Inches(0.98), Inches(5.96), Inches(9.0), Inches(0.68), LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
add_text(s, Inches(1.22), Inches(5.96), Inches(8.6), Inches(0.68),
         [[("Palavras-chave:  ", SZ_BODY_SM, TEAL, True),
           ("TSP · Quantum Annealing · Otimização Combinatória "
            "· Logística", SZ_BODY_SM, INK)]],
         anchor=MSO_ANCHOR.MIDDLE)
add_pic_fit(s, LOGO_CIMATEC, Inches(10.15), Inches(5.78), Inches(2.4),
            Inches(1.0), align="right")
add_text(s, Inches(0.98), Inches(6.86), Inches(11.0), Inches(0.4),
         [[("XI Seminário de Avaliação de Pesquisa Científica "
            "e Tecnológica — SAPCT 2026", SZ_FOOTER, GRAY, False, True)]])
notes(s, "Abertura (~20s). Apresentar titulo, autores e a natureza do "
         "trabalho: revisao exploratoria da literatura.")

# ----------------------------------------------------------------------
# SLIDE 3 — Resumo
# ----------------------------------------------------------------------
s = content_slide(prs, "Resumo", "Síntese do trabalho", 3)
add_bullets(s, Inches(0.66), Inches(1.98), Inches(7.20), Inches(4.0), [
    ("O TSP é um problema de otimização combinatória NP-Hard, de "
     "grande relevância para a logística.",),
    ("O Quantum Annealing (QA) destaca-se como abordagem promissora "
     "frente à complexidade computacional do problema.",),
    ("Revisão exploratória da literatura: 26 strings de busca na base "
     "Lens.org resultaram em 3.696 artigos únicos.",),
    ("O QA é o paradigma dominante, com destaque para a "
     "plataforma D-Wave.",),
    ("Há lacunas em abordagens gate-based e em aprendizado de "
     "máquina quântico para roteamento logístico.",),
], size=SZ_BODY, gap=12)
# indicadores
stat_card(s, Inches(8.10), Inches(1.98), Inches(4.57), Inches(1.42),
          "3.696", "Artigos únicos analisados (após deduplicação)")
stat_card(s, Inches(8.10), Inches(3.54), Inches(4.57), Inches(1.42),
          "26", "Strings de busca executadas na base Lens.org")
# barra de palavras-chave
add_rect(s, Inches(0.66), Inches(5.30), Inches(12.01), Inches(0.78), NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
add_text(s, Inches(1.0), Inches(5.30), Inches(11.4), Inches(0.78),
         [[("Palavras-chave:   ", SZ_BODY, STEEL, True),
           ("Problema do Caixeiro Viajante  ·  Quantum Annealing  ·  "
            "Otimização Combinatória  ·  Logística", SZ_BODY, WHITE)]],
         anchor=MSO_ANCHOR.MIDDLE)
notes(s, "Resumo (~50s). Dar a visao geral: problema (TSP NP-Hard), "
         "tecnologia (QA), metodo (26 strings -> 3.696 artigos) e o "
         "achado central (QA dominante; lacunas gate-based e QML).")

# ----------------------------------------------------------------------
# SLIDE 4 — Introducao
# ----------------------------------------------------------------------
s = content_slide(prs, "Introdução", "Contexto e objetivo", 4)
# coluna 1 — o problema
add_rect(s, Inches(0.66), Inches(1.92), Inches(5.97), Inches(2.92), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
add_rect(s, Inches(0.66), Inches(1.92), Inches(5.97), Inches(0.08), NAVY)
add_text(s, Inches(0.94), Inches(2.10), Inches(5.4), Inches(0.4),
         [[("O problema", SZ_SUBHEAD, NAVY, True)]])
add_bullets(s, Inches(0.94), Inches(2.58), Inches(5.42), Inches(2.2), [
    ("TSP: rota de menor custo que visita cada cidade uma vez e "
     "retorna à origem — classificado como NP-Hard.",),
    ("Relevância logística direta: otimizar rotas reduz custos "
     "operacionais e eleva a eficiência.",),
    ("TSPTW: variante com janelas de tempo, próxima de cenários "
     "reais de distribuição e roteamento.",),
], size=SZ_BODY, gap=8)
# coluna 2 — a tecnologia
add_rect(s, Inches(6.72), Inches(1.92), Inches(5.95), Inches(2.92), NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
add_rect(s, Inches(6.72), Inches(1.92), Inches(5.95), Inches(0.08), TEAL)
add_text(s, Inches(7.00), Inches(2.10), Inches(5.4), Inches(0.4),
         [[("A tecnologia", SZ_SUBHEAD, WHITE, True)]])
add_bullets(s, Inches(7.00), Inches(2.58), Inches(5.40), Inches(2.2), [
    ("Computação quântica explora o espaço de soluções de forma "
     "mais eficiente que métodos clássicos.",),
    ("Quantum Annealing: tunelamento quântico para o estado de "
     "menor energia; formulação QUBO; hardware D-Wave.",),
    ("Salehi, Glos e Miszczak (2022): três formulações QUBO "
     "para o TSPTW.",),
], size=SZ_BODY, gap=8, color=PALE, lead_color=STEEL)
# objetivo
add_rect(s, Inches(0.66), Inches(5.00), Inches(12.01), Inches(1.42), LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
add_rect(s, Inches(0.66), Inches(5.00), Inches(0.10), Inches(1.42), TEAL)
add_text(s, Inches(1.04), Inches(5.14), Inches(11.4), Inches(0.36),
         [[("OBJETIVO", SZ_KICKER, TEAL, True)]])
add_text(s, Inches(1.04), Inches(5.46), Inches(11.4), Inches(0.88),
         [[("Realizar uma revisão exploratória da literatura sobre a "
            "aplicação de algoritmos quânticos, com ênfase em Quantum "
            "Annealing, ao TSP e suas variantes em logística — mapeando o "
            "volume de produção, os algoritmos mais utilizados e as "
            "lacunas de pesquisa.", SZ_BODY, INK)]], line_spacing=1.18)
notes(s, "Introducao (~90s). Esquerda: por que o TSP importa (NP-Hard, "
         "logistica, TSPTW). Direita: por que a computacao quantica/QA e "
         "uma alternativa (QUBO, D-Wave, Salehi 2022). Fechar lendo o "
         "objetivo da revisao.")

# ----------------------------------------------------------------------
# SLIDE 5 — Metodologia
# ----------------------------------------------------------------------
s = content_slide(prs, "Metodologia", "Como a literatura foi mapeada", 5)
add_text(s, Inches(0.66), Inches(1.90), Inches(6.1), Inches(0.4),
         [[("Pesquisa bibliográfica exploratória em três eixos",
            SZ_SUBHEAD, BLUE, True)]])
eixos = [
    ("1 · O problema", "TSP e variantes — VRP, CVRP, Job Shop "
     "Scheduling, Bin Packing."),
    ("2 · A tecnologia", "Paradigmas quânticos — QA, QAOA, VQE, "
     "Grover, QUBO/Ising."),
    ("3 · A aplicação", "Logística, cadeia de suprimentos e "
     "otimização de rotas."),
]
for i, (t, d) in enumerate(eixos):
    tp = Inches(2.36) + i * Inches(0.86)
    add_rect(s, Inches(0.66), tp, Inches(6.05), Inches(0.74), CARD,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    add_rect(s, Inches(0.66), tp, Inches(0.08), Inches(0.74), TEAL)
    add_text(s, Inches(0.92), tp, Inches(1.92), Inches(0.74),
             [[(t, SZ_BODY, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.0)
    add_text(s, Inches(2.92), tp, Inches(3.66), Inches(0.74),
             [[(d, SZ_BODY_SM, GRAY)]], anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.06)
add_bullets(s, Inches(0.66), Inches(5.12), Inches(6.05), Inches(1.7), [
    ("26 strings de busca (operador AND) na base Lens.org; "
     "prioridades 8 Alta · 10 Média · 8 Baixa.",),
    ("Duas fases: 16 strings iniciais + 10 derivadas de "
     "Phillipson (2025).",),
    ("Deduplicação em duas etapas: por DOI e por título "
     "normalizado.",),
], size=SZ_BODY, gap=8)
add_rect(s, Inches(6.92), Inches(1.86), Inches(5.75), Inches(4.55), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_pic_fit(s, os.path.join(IMG, "06_strings_busca_01_volume_strings.png"),
            Inches(7.00), Inches(1.96), Inches(5.59), Inches(4.05))
caption(s, Inches(6.92), Inches(6.02), Inches(5.75),
        "Volume de artigos únicos recuperados por string de busca.")
notes(s, "Metodologia (~90s). Tres eixos tematicos cruzados em 26 strings "
         "no Lens.org, priorizadas e aplicadas em duas fases. Deduplicacao "
         "por DOI e titulo. O grafico mostra a contribuicao de cada string.")

# ----------------------------------------------------------------------
# SLIDE 6 — Resultados e Discussao I — Panorama
# ----------------------------------------------------------------------
s = content_slide(prs, "Resultados e Discussão · 1",
                   "Panorama: um campo em expansão", 6)
cards = [
    ("3.696", "Artigos únicos (de 5.320 brutos — 30,5% duplicatas)"),
    ("61,7%", "Publicados a partir de 2020 — pico em 2025"),
    ("49", "Países representados na produção científica"),
    ("82,8%", "Artigos em acesso aberto (Open Access)"),
]
cw = Inches(2.92)
gx = Inches(0.135)
for i, (b, lab) in enumerate(cards):
    l = Inches(0.66) + i * (cw + gx)
    stat_card(s, l, Inches(1.94), cw, Inches(1.36), b, lab)
add_rect(s, Inches(0.66), Inches(3.48), Inches(12.01), Inches(2.62), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_pic_fit(s, os.path.join(IMG, "01_visao_geral_01_publicacoes_tempo.png"),
            Inches(0.82), Inches(3.56), Inches(11.69), Inches(2.34))
caption(s, Inches(0.66), Inches(6.04), Inches(12.01),
        "Publicações ao longo do tempo — forte aceleração a partir de "
        "2020. Produção globalmente distribuída e com forte adesão à "
        "ciência aberta, o que favorece a reprodutibilidade.")
notes(s, "Resultados 1 (~70s). De 5.320 registros brutos para 3.696 unicos. "
         "61,7% pos-2020, pico em 2025: campo em franca expansao. 49 paises "
         "e 82,8% de acesso aberto reforcam alcance global e "
         "reprodutibilidade.")

# ----------------------------------------------------------------------
# SLIDE 7 — Resultados e Discussao II — Paradigmas e algoritmos
# ----------------------------------------------------------------------
s = content_slide(prs, "Resultados e Discussão · 2",
                   "Paradigma dominante e taxonomia", 7)
add_bullets(s, Inches(0.66), Inches(2.04), Inches(6.05), Inches(4.2), [
    ("O Quantum Annealing é o paradigma dominante: mais de 1.200 "
     "artigos nas strings diretamente relacionadas.",),
    ("Hardware de referência — D-Wave Advantage: 5.640 qubits e "
     "conectividade Pegasus.",),
    ("Algoritmos gate-based (QAOA, VQE) somam ≈ 150 artigos "
     "aplicados ao TSP.",),
    ("Mais de 40 algoritmos catalogados em 6 categorias; predominam "
     "as abordagens híbridas (quântico + clássico).",),
], size=SZ_BODY, gap=13)
add_rect(s, Inches(6.86), Inches(1.92), Inches(5.81), Inches(2.22), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
add_pic_fit(s, os.path.join(IMG, "07_algoritmos_01_paradigma.png"),
            Inches(6.94), Inches(1.98), Inches(5.65), Inches(1.96))
add_rect(s, Inches(6.86), Inches(4.26), Inches(5.81), Inches(2.22), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
add_pic_fit(s, os.path.join(IMG, "07_algoritmos_02_abordagem.png"),
            Inches(6.94), Inches(4.32), Inches(5.65), Inches(1.96))
caption(s, Inches(6.86), Inches(6.40), Inches(5.81),
        "Amostra de artigos classificados — paradigma e abordagem.")
notes(s, "Resultados 2 (~80s). QA domina (>1.200 artigos; D-Wave "
         "Advantage). Gate-based (QAOA/VQE) ainda incipiente (~150). "
         "Taxonomia: 40+ algoritmos em 6 categorias; a maioria das "
         "implementacoes e hibrida.")

# ----------------------------------------------------------------------
# SLIDE 8 — Resultados e Discussao III — Lacunas
# ----------------------------------------------------------------------
s = content_slide(prs, "Resultados e Discussão · 3",
                   "Lacunas de pesquisa identificadas", 8)
add_text(s, Inches(0.66), Inches(1.90), Inches(11.9), Inches(0.4),
         [[("Quatro lacunas se destacam e apontam direções para "
            "trabalhos futuros:", SZ_SUBHEAD, BLUE, True)]])
lac = [
    ("Gate-based", "Escassas",
     "Faltam soluções QAOA/VQE com escalabilidade demonstrada "
     "para o TSP."),
    ("QRL para roteamento", "21 artigos",
     "Quantum Reinforcement Learning para roteamento — área "
     "emergente e pouco explorada."),
    ("CVRP com QA", "18 artigos",
     "Aplicação logística direta do Quantum Annealing ao CVRP "
     "ainda pouco investigada."),
    ("Dados reais", "Limitada",
     "Validação com dados logísticos reais — predominam "
     "benchmarks acadêmicos."),
]
cw = Inches(2.92)
gx = Inches(0.135)
for i, (t, big, d) in enumerate(lac):
    l = Inches(0.66) + i * (cw + gx)
    add_rect(s, l, Inches(2.46), cw, Inches(3.42), CARD,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    add_rect(s, l, Inches(2.46), cw, Inches(0.92), NAVY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.20)
    add_rect(s, l, Inches(2.96), cw, Inches(0.42), NAVY)
    add_text(s, l, Inches(2.46), cw, Inches(0.92),
             [[(t, SZ_BODY, WHITE, True)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, l + Inches(0.18), Inches(3.52), cw - Inches(0.36), Inches(0.6),
             [[(big, SZ_STAT, TEAL, True)]], align=PP_ALIGN.CENTER)
    add_text(s, l + Inches(0.26), Inches(4.24), cw - Inches(0.52),
             Inches(1.5), [[(d, SZ_BODY, INK)]], line_spacing=1.18,
             align=PP_ALIGN.CENTER)
add_text(s, Inches(0.66), Inches(6.06), Inches(12.0), Inches(0.4),
         [[("Persiste o desafio de demonstrar escalabilidade frente às "
            "limitações do hardware NISQ.", SZ_BODY_SM, GRAY, False, True)]],
         align=PP_ALIGN.CENTER)
notes(s, "Resultados 3 (~60s). Quatro lacunas: gate-based escalavel, QRL "
         "para roteamento (21), CVRP com QA (18) e validacao com dados "
         "reais. Sao oportunidades concretas de contribuicao.")

# ----------------------------------------------------------------------
# SLIDE 9 — Consideracoes Finais
# ----------------------------------------------------------------------
s = content_slide(prs, "Considerações Finais",
                   "Conclusões e próximas etapas", 9)
add_rect(s, Inches(0.66), Inches(1.96), Inches(5.97), Inches(4.46), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_rect(s, Inches(0.66), Inches(1.96), Inches(5.97), Inches(0.08), NAVY)
add_text(s, Inches(0.96), Inches(2.16), Inches(5.4), Inches(0.4),
         [[("O que a revisão mostrou", SZ_SUBHEAD, NAVY, True)]])
add_bullets(s, Inches(0.96), Inches(2.66), Inches(5.42), Inches(3.6), [
    ("A aplicação de computação quântica ao TSP é um campo em "
     "rápida expansão, com base robusta de 3.696 artigos.",),
    ("O Quantum Annealing consolidou-se como o paradigma mais "
     "maduro e acessível para otimização combinatória.",),
    ("Algoritmos variacionais e de aprendizado de máquina quântico "
     "são fronteiras ainda pouco exploradas.",),
    ("As lacunas identificadas são oportunidades concretas de "
     "contribuição científica.",),
], size=SZ_BODY, gap=12)
add_rect(s, Inches(6.72), Inches(1.96), Inches(5.95), Inches(4.46), NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_rect(s, Inches(6.72), Inches(1.96), Inches(5.95), Inches(0.08), TEAL)
add_text(s, Inches(7.02), Inches(2.16), Inches(5.4), Inches(0.4),
         [[("Próximas etapas da pesquisa", SZ_SUBHEAD, WHITE, True)]])
steps = [
    "Fundamentação teórica aprofundada do tema.",
    "Seleção e reprodução de algoritmos em simuladores quânticos.",
    "Comparação experimental entre abordagens clássicas e quânticas "
    "para instâncias do TSP em contexto logístico.",
]
for i, t in enumerate(steps):
    tp = Inches(2.80) + i * Inches(1.18)
    add_rect(s, Inches(7.02), tp, Inches(0.56), Inches(0.56), TEAL,
             shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(7.02), tp, Inches(0.56), Inches(0.56),
             [[(str(i + 1), SZ_SUBHEAD, WHITE, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.76), tp - Inches(0.06), Inches(4.7), Inches(1.1),
             [[(t, SZ_BODY, PALE)]], anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.18)
notes(s, "Consideracoes finais (~60s). Sintetizar: campo em expansao, QA "
         "maduro, variacionais/QML como fronteiras. Apresentar as tres "
         "proximas etapas da pesquisa de mestrado.")

# ----------------------------------------------------------------------
# SLIDE 10 — Agradecimentos
# ----------------------------------------------------------------------
s = content_slide(prs, "Agradecimentos", "Agradecimentos", 10)
add_rect(s, Inches(0.66), Inches(2.10), Inches(12.01), Inches(2.20), LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
add_rect(s, Inches(0.66), Inches(2.10), Inches(0.10), Inches(2.20), TEAL)
add_text(s, Inches(1.10), Inches(2.42), Inches(11.2), Inches(1.6),
         [[("Este trabalho foi parcialmente financiado pelo projeto QuIIN "
            "FCRH Mestrado TQ, suportado pelo QuIIN — Quantum Industrial "
            "Innovation, Centro de Competência EMBRAPII CIMATEC em "
            "Tecnologias Quânticas, com recursos do PPI IoT/Manufatura 4.0 "
            "do MCTI, por meio do Termo de Cooperação 053/2023 firmado com "
            "a EMBRAPII. Apoio adicional do CNPq — Conselho Nacional de "
            "Desenvolvimento Científico e Tecnológico.", SZ_BODY, INK)]],
         line_spacing=1.28)
add_text(s, Inches(0.66), Inches(4.70), Inches(11.0), Inches(0.4),
         [[("Apoio institucional e fomento", SZ_KICKER, TEAL, True)]])
entes = ["QuIIN", "EMBRAPII CIMATEC", "MCTI · PPI IoT / Manufatura 4.0",
         "CNPq"]
cw = Inches(2.92)
gx = Inches(0.135)
for i, e in enumerate(entes):
    l = Inches(0.66) + i * (cw + gx)
    chip(s, l, Inches(5.10), cw, e)
notes(s, "Agradecimentos (~20s). Reconhecer o fomento: QuIIN, EMBRAPII "
         "CIMATEC, MCTI e CNPq.")

# ----------------------------------------------------------------------
# SLIDE 11 — Referencias
# ----------------------------------------------------------------------
s = content_slide(prs, "Referências", "Referências", 11)
refs = [
    "APPLEGATE, D. L. et al. The Traveling Salesman Problem: a "
    "computational study. Princeton: Princeton University Press, 2006.",
    "KADOWAKI, T.; NISHIMORI, H. Quantum annealing in the Transverse "
    "Ising model. Physical Review E, v. 58, n. 5, p. 5355-5363, 1998.",
    "SALEHI, Ö.; GLOS, A.; MISZCZAK, J. A. Unconstrained binary models "
    "of the travelling salesman problem variants for quantum "
    "optimization. Quantum Information Processing, v. 21, n. 67, 2022.",
    "PHILLIPSON, F. Quantum Computing in Logistics and Supply Chain "
    "Management: an overview. Maastricht University / TNO, 2025.",
    "FARHI, E.; GOLDSTONE, J.; GUTMANN, S. A quantum approximate "
    "optimization algorithm. arXiv preprint arXiv:1411.4028, 2014.",
]
add_rect(s, Inches(0.66), Inches(1.98), Inches(12.01), Inches(4.42), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_rect(s, Inches(0.66), Inches(1.98), Inches(12.01), Inches(0.08), NAVY)
tb = s.shapes.add_textbox(Inches(1.04), Inches(2.34), Inches(11.3),
                          Inches(3.8))
tf = tb.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
for i, r in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.line_spacing = 1.2
    p.space_after = Pt(14)
    num = p.add_run()
    num.text = "%d.   " % (i + 1)
    _style_run(num, SZ_BODY, TEAL, True)
    body = p.add_run()
    body.text = r
    _style_run(body, SZ_BODY, INK)
notes(s, "Referencias (~20s). Destacar as obras-chave: Kadowaki & Nishimori "
         "(fundamento do QA), Salehi et al. (QUBO para TSPTW) e Phillipson "
         "(base da segunda fase de busca).")

# ----------------------------------------------------------------------
# SLIDE 12 — Encerramento
# ----------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[0])
clear_placeholders(s)
set_bg(s, NAVY)
add_rect(s, 0, Inches(6.86), EMU_W, Inches(0.64), TEAL)
add_pic_fit(s, LOGO_SAPCT, Inches(0.95), Inches(0.78), Inches(2.7),
            Inches(1.45), align="left")
add_text(s, Inches(0.98), Inches(2.66), Inches(10.0), Inches(1.3),
         [[("Obrigado!", 52, WHITE, True)]])
add_rect(s, Inches(1.02), Inches(3.86), Inches(1.15), Pt(4), TEAL)
add_text(s, Inches(0.98), Inches(4.06), Inches(11.2), Inches(0.95),
         [[("Aplicação de Quantum Annealing ao Problema do "
            "Caixeiro Viajante:", SZ_SUBHEAD, STEEL, True)],
          [("revisão exploratória da literatura em logística",
            SZ_BODY, RGBColor(0xB9, 0xC6, 0xD6), False, True)]],
         line_spacing=1.12, space_after=5)
add_text(s, Inches(0.98), Inches(5.22), Inches(11.2), Inches(0.5),
         [[("Renan Cardoso dos Santos", SZ_BODY, WHITE, True),
           ("   ·   Valéria Loureiro da Silva   ·   "
            "Anderson Rafael Correia B. da Silva", SZ_BODY_SM,
            RGBColor(0xB9, 0xC6, 0xD6))]])
add_text(s, Inches(0.98), Inches(5.72), Inches(11.0), Inches(0.5),
         [[("✉  ", SZ_BODY, TEAL, True),
           ("renan.santos@fbter.org.br", SZ_BODY, STEEL),
           ("        Universidade SENAI CIMATEC · Salvador-BA",
            SZ_BODY_SM, RGBColor(0x8F, 0x9F, 0xB3))]])
add_text(s, Inches(0.98), Inches(6.97), Inches(9.0), Inches(0.42),
         [[("XI SAPCT — SENAI CIMATEC · 28 e 29 de Maio de 2026",
            SZ_FOOTER, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, Inches(10.62), Inches(6.92), Inches(2.10), Inches(0.50), WHITE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
add_pic_fit(s, LOGO_CIMATEC, Inches(10.74), Inches(6.96), Inches(1.86),
            Inches(0.42), align="center")
notes(s, "Encerramento (~10s). Agradecer e abrir para perguntas.")

# ----------------------------------------------------------------------
prs.save(OUTPUT)
print("OK ->", OUTPUT)
print("Total de slides:", len(prs.slides._sldIdLst))
