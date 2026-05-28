# -*- coding: utf-8 -*-
"""
Aplica as 6 recomendacoes de melhoria na apresentacao v2.1 -> v2.2.

Origem do trabalho: ambiguidade entre dois universos de dados que aparecem
lado a lado nos slides 9-12 sem sinalizacao visual:
  - Base A: corpus bibliometrico (3.696 artigos do Lens.org)
  - Base B: amostra curada Phillipson 2025 (129 trabalhos catalogados)

Recomendacoes aplicadas (R1-R6):
  R1. Slide 8 (Metodologia)  -> nota explicativa das duas bases
  R2. Slides 9-12            -> badges de fonte (Base A / Base B)
  R3. Slide 10 (paradigma)   -> KPI duplicado em duas colunas:
                                (esq) 1.117 menções QA no corpus  (Base A)
                                (dir) 53,5% dos 129 catalogados   (Base B)
  R4. Slide 11 (taxonomia)   -> headline e caption mencionam "(n=129, Phillipson 2025)"
  R5. Slide 12 (lacunas)     -> badge Base A (numeros sao do corpus bibliometrico)

Execucao:
    python src/aplicar_melhorias_apresentacao.py

Entrada:  docs/Apresentacao_SAPCT_CIMATEC_2026_TSP_Quantum_Annealing_v2.1.pptx
Saida:    docs/Apresentacao_SAPCT_CIMATEC_2026_TSP_Quantum_Annealing_v2.2.pptx
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# CAMINHOS
# ============================================================
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PATH_IN  = os.path.join(ROOT, "docs", "Apresentacao_SAPCT_CIMATEC_2026_TSP_Quantum_Annealing_v2.1.pptx")
PATH_OUT = os.path.join(ROOT, "docs", "Apresentacao_SAPCT_CIMATEC_2026_TSP_Quantum_Annealing_v2.2.pptx")


# ============================================================
# CORES (alinhadas com a paleta hibrida do dashboard)
# ============================================================
COR_BASE_A     = RGBColor(0x3C, 0x60, 0xA7)   # Pantone 293 — corpus bibliometrico
COR_BASE_B     = RGBColor(0xD9, 0x77, 0x57)   # Terracota   — amostra curada Phillipson
COR_BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
COR_AZUL_KPI   = RGBColor(0x0E, 0x28, 0x41)   # azul escuro do KPI original
COR_CINZA_BG   = RGBColor(0xF5, 0xF6, 0xF9)   # cinza muito claro para fundo do KPI


# ============================================================
# HELPERS
# ============================================================
def add_badge(slide, x_in, y_in, w_in, h_in, cor_bg, texto, font_size=10):
    """Caixa arredondada colorida com texto branco — indica fonte de dados."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor_bg
    shape.line.fill.background()  # sem borda
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = COR_BRANCO
    return shape


def add_textbox(slide, x_in, y_in, w_in, h_in, texto,
                font_size=11, bold=False, color=None, align=PP_ALIGN.LEFT):
    """Caixa de texto simples (sem fundo, sem borda)."""
    tb = slide.shapes.add_textbox(
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return tb


def find_shape_by_text(slide, fragmento):
    """Retorna o primeiro shape cujo text_frame contem o fragmento."""
    for shape in slide.shapes:
        if shape.has_text_frame and fragmento in shape.text_frame.text:
            return shape
    return None


def replace_text_preservando_formato(shape, novo_texto):
    """Substitui o texto de um shape mantendo o formato do primeiro run.

    PowerPoint pode quebrar um texto em multiplos runs (ex: "+1.200" em "+" + "1.200").
    Esta funcao concentra o texto no primeiro run e remove os demais.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    p0 = tf.paragraphs[0]
    if p0.runs:
        # Usa o primeiro run (mantem formato), substitui texto
        p0.runs[0].text = novo_texto
        # Remove runs adicionais do primeiro paragraph
        for run in list(p0.runs[1:]):
            run._r.getparent().remove(run._r)
    else:
        # Cria run novo se nao existir
        run = p0.add_run()
        run.text = novo_texto
    # Remove paragraphs extras
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)


def criar_kpi_box(slide, x_in, y_in, w_in, h_in,
                  valor, label, cor_barra):
    """Reproduz visualmente o padrao de KPI do slide 10 (fundo + barra lateral + valor + label)."""
    # Fundo (rounded rectangle cinza claro)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COR_CINZA_BG
    bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    bg.line.width = Pt(0.5)

    # Barra lateral colorida (decorativa)
    barra = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_in), Inches(y_in), Inches(0.07), Inches(h_in),
    )
    barra.fill.solid()
    barra.fill.fore_color.rgb = cor_barra
    barra.line.fill.background()

    # Valor (numero grande)
    add_textbox(slide,
        x_in=x_in + 0.18, y_in=y_in + 0.16, w_in=w_in - 0.30, h_in=0.40,
        texto=valor,
        font_size=23, bold=True, color=COR_AZUL_KPI,
    )

    # Label (texto descritivo)
    add_textbox(slide,
        x_in=x_in + 0.18, y_in=y_in + 0.63, w_in=w_in - 0.30, h_in=0.70,
        texto=label,
        font_size=10, bold=False, color=COR_AZUL_KPI,
    )


# ============================================================
# APLICAR RECOMENDACOES
# ============================================================
def aplicar_r1_slide8_nota_bases(slide):
    """R1: nota explicativa das duas bases no rodape do slide de Metodologia."""
    add_textbox(slide,
        x_in=0.66, y_in=6.65, w_in=12.0, h_in=0.34,
        texto="(A) Corpus bibliométrico (n=3.696, Lens.org)   ·   "
              "(B) Amostra Phillipson 2025 (n=129) — caracterizam o campo em escalas diferentes",
        font_size=10, bold=True, color=COR_AZUL_KPI,
    )


def aplicar_r2_slide9_badge_base_a(slide):
    """R2: badge Base A no slide 9 (Resultados-1 — panorama bibliometrico)."""
    add_badge(slide,
        x_in=8.30, y_in=1.05, w_in=4.30, h_in=0.40,
        cor_bg=COR_BASE_A,
        texto="(A)  Corpus bibliométrico  ·  n=3.696",
        font_size=10,
    )


def aplicar_r3_slide10_duas_colunas(slide):
    """R3: reformula slide 10 — KPI Base A (esquerda) + KPI Base B (direita)."""
    # 3a) Corrige o numero "+1.200" -> "1.117" (valor verificado no corpus)
    sh_valor = find_shape_by_text(slide, "+1.200")
    if sh_valor is None:
        print("  [AVISO] shape '+1.200' nao encontrado no slide 10")
    else:
        replace_text_preservando_formato(sh_valor, "1.117")

    # 3b) Atualiza label do KPI esquerdo (deixa explicita a origem na Base A)
    sh_label = find_shape_by_text(slide, "Artigos em Quantum Annealing")
    if sh_label is None:
        print("  [AVISO] label do KPI esquerdo nao encontrado")
    else:
        replace_text_preservando_formato(
            sh_label,
            "Menções a QA no corpus bibliométrico (6 das 26 strings)",
        )

    # 3c) Badge Base A acima do KPI esquerdo
    add_badge(slide,
        x_in=0.62, y_in=1.45, w_in=2.86, h_in=0.28,
        cor_bg=COR_BASE_A,
        texto="(A)  Base bibliométrica  ·  n=3.696",
        font_size=9,
    )

    # 3d) Novo KPI Base B ao lado direito do existente
    # KPI atual: x=0.62, y=1.80, 2.86x1.41
    # Novo:      x=3.62, y=1.80, 2.86x1.41 (mesmo tamanho, dx=+3.0)
    criar_kpi_box(slide,
        x_in=3.62, y_in=1.80, w_in=2.86, h_in=1.41,
        valor="53,5%",
        label="Dos 129 trabalhos catalogados em Phillipson (2025) usam QA",
        cor_barra=COR_BASE_B,
    )

    # 3e) Badge Base B acima do KPI direito
    add_badge(slide,
        x_in=3.62, y_in=1.45, w_in=2.86, h_in=0.28,
        cor_bg=COR_BASE_B,
        texto="(B)  Amostra Phillipson 2025  ·  n=129",
        font_size=9,
    )

    # 3f) Atualiza caption das imagens (graficos sao da Base B)
    sh_cap = find_shape_by_text(slide, "Amostra de artigos classificados")
    if sh_cap is None:
        print("  [AVISO] caption das imagens nao encontrada")
    else:
        replace_text_preservando_formato(
            sh_cap,
            "Distribuição por paradigma (acima) e algoritmos mais frequentes (abaixo) "
            "— amostra de 129 trabalhos classificados em Phillipson (2025).",
        )


def aplicar_r4_slide11_headline_caption(slide):
    """R4: explicita 'n=129, Phillipson 2025' na headline e caption do slide 11."""
    # 4a) Headline: "Mais de 40 algoritmos catalogados" -> "40+ algoritmos identificados na amostra Phillipson (2025)"
    sh_head = find_shape_by_text(slide, "Mais de 40 algoritmos catalogados")
    if sh_head is None:
        print("  [AVISO] headline do slide 11 nao encontrada")
    else:
        replace_text_preservando_formato(
            sh_head,
            "40+ algoritmos identificados na amostra Phillipson (2025, n=129) → 6 categorias principais:",
        )

    # 4b) Caption do grafico de abordagem
    sh_cap = find_shape_by_text(slide, "Distribuição por abordagem na amostra")
    if sh_cap is None:
        print("  [AVISO] caption de abordagem nao encontrada")
    else:
        replace_text_preservando_formato(
            sh_cap,
            "Distribuição por abordagem na amostra Phillipson (n=129, 2025) "
            "— predomínio de soluções híbridas.",
        )

    # 4c) Badge Base B no canto superior direito
    add_badge(slide,
        x_in=8.30, y_in=1.05, w_in=4.30, h_in=0.40,
        cor_bg=COR_BASE_B,
        texto="(B)  Amostra Phillipson 2025  ·  n=129",
        font_size=10,
    )


def aplicar_r5_slide12_badge_base_a(slide):
    """R5: badge Base A no slide 12 (numeros 21 e 18 vem do corpus, nao da curada)."""
    add_badge(slide,
        x_in=8.30, y_in=1.05, w_in=4.30, h_in=0.40,
        cor_bg=COR_BASE_A,
        texto="(A)  Contagens das strings de busca",
        font_size=10,
    )


# ============================================================
# MAIN
# ============================================================
def main():
    if not os.path.exists(PATH_IN):
        print(f"[ERRO] Arquivo de entrada nao encontrado: {PATH_IN}")
        sys.exit(1)

    print(f"Lendo:  {os.path.basename(PATH_IN)}")
    prs = Presentation(PATH_IN)
    slides = list(prs.slides)

    if len(slides) < 12:
        print(f"[ERRO] Esperado >= 12 slides, encontrado {len(slides)}")
        sys.exit(1)

    s8, s9, s10, s11, s12 = slides[7], slides[8], slides[9], slides[10], slides[11]

    print("\nAplicando R1 — nota das duas bases (slide 8)")
    aplicar_r1_slide8_nota_bases(s8)

    print("Aplicando R2 — badge Base A (slide 9)")
    aplicar_r2_slide9_badge_base_a(s9)

    print("Aplicando R3 — reformulacao do slide 10 (duas colunas)")
    aplicar_r3_slide10_duas_colunas(s10)

    print("Aplicando R4 — headline + caption + badge B (slide 11)")
    aplicar_r4_slide11_headline_caption(s11)

    print("Aplicando R5 — badge Base A (slide 12)")
    aplicar_r5_slide12_badge_base_a(s12)

    print(f"\nSalvando: {os.path.basename(PATH_OUT)}")
    prs.save(PATH_OUT)
    print("Concluido.")


if __name__ == "__main__":
    main()
