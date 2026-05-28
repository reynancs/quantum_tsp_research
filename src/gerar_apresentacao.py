# -*- coding: utf-8 -*-
"""
Gera a apresentacao oral (.pptx) do artigo de revisao bibliografica
"Aplicacao de Quantum Annealing ao Problema do Caixeiro Viajante" para o
XI SAPCT - SENAI CIMATEC 2026, usando o template padrao do evento.

Conteudo extraido exclusivamente de:
  docs/resumo_final_expandido_SAPCT_CIMATEC_2026.pdf
Imagens: docs/images/ (graficos do levantamento bibliometrico).

Execucao:  python src/gerar_apresentacao.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

# ----------------------------------------------------------------------
# Caminhos
# ----------------------------------------------------------------------
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS = os.path.join(ROOT, "docs")
IMG = os.path.join(DOCS, "images")
MEDIA = os.path.join(DOCS, "tmpl_unpacked", "ppt", "media")
TEMPLATE = os.path.join(DOCS, "Template_Apresentacao_Oral_SAPCTSLIDES2026B.pptx")
OUTPUT = os.path.join(DOCS, "Apresentacao_SAPCT_CIMATEC_2026_TSP_Quantum_Annealing.pptx")

LOGO_SAPCT = os.path.join(MEDIA, "image4.png")    # marca SAPCT (cerebro)
LOGO_CIMATEC = os.path.join(MEDIA, "image5.png")  # marca SENAI CIMATEC

# ----------------------------------------------------------------------
# Paleta de cores (alinhada ao tema do template e aos graficos)
# ----------------------------------------------------------------------
NAVY = RGBColor(0x0E, 0x28, 0x41)   # azul-marinho - titulos / destaque
BLUE = RGBColor(0x2E, 0x5C, 0x8A)   # azul medio - subtitulos
TEAL = RGBColor(0x15, 0x60, 0x82)   # accent1 do tema - reguas/realces
STEEL = RGBColor(0x5B, 0x7F, 0xB5)  # azul claro dos graficos
LIGHT = RGBColor(0xEC, 0xF1, 0xF6)  # painel claro
CARD = RGBColor(0xF4, 0xF6, 0xF9)   # cartao claro
INK = RGBColor(0x2B, 0x2B, 0x2B)    # texto corrido
GRAY = RGBColor(0x6B, 0x72, 0x7B)   # texto secundario
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Arial"
FOOTER_TXT = "XI SAPCT  ·  SENAI CIMATEC  ·  Salvador-BA  ·  28 e 29 de Maio de 2026"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

# ----------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------

def set_bg(slide, color=WHITE):
    """Define cor solida de fundo do slide."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def clear_placeholders(slide):
    """Remove placeholders herdados do layout (data/rodape/numero)."""
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def _style_run(run, size, color, bold=False, italic=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    # garante a fonte tambem para caracteres latinos/cs
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def add_text(slide, left, top, width, height, runs_spec, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=4, wrap=True):
    """Cria uma caixa de texto.

    runs_spec: lista de paragrafos; cada paragrafo e' uma lista de tuplas
    (texto, size, color, bold, italic).
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs_spec):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for spec in para:
            text, size, color = spec[0], spec[1], spec[2]
            bold = spec[3] if len(spec) > 3 else False
            italic = spec[4] if len(spec) > 4 else False
            r = p.add_run()
            r.text = text
            _style_run(r, size, color, bold, italic)
    return tb


def add_bullets(slide, left, top, width, height, items, size=15,
                color=INK, gap=7, lead_color=TEAL, line_spacing=1.08):
    """Lista de topicos. items: (texto, nivel, bold?)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        text = it[0]
        level = it[1] if len(it) > 1 else 0
        bold = it[2] if len(it) > 2 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        if level == 0:
            mark = p.add_run()
            mark.text = "▪  "
            _style_run(mark, size, lead_color, True)
            r = p.add_run()
            r.text = text
            _style_run(r, size, color, bold)
        else:
            p.level = 1
            mark = p.add_run()
            mark.text = "      –  "
            _style_run(mark, size - 1, STEEL, False)
            r = p.add_run()
            r.text = text
            _style_run(r, size - 1, GRAY, False)
    return tb


def add_pic_fit(slide, path, box_l, box_t, box_w, box_h, align="center"):
    """Insere imagem ajustada (preservando proporcao) dentro de uma caixa."""
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    t = box_t + (box_h - h) // 2
    if align == "center":
        l = box_l + (box_w - w) // 2
    elif align == "left":
        l = box_l
    else:
        l = box_l + (box_w - w)
    return slide.shapes.add_picture(path, l, t, w, h)


def add_rect(slide, left, top, width, height, fill, shape=MSO_SHAPE.RECTANGLE,
             line=None, line_w=1.0, radius=None, shadow=False):
    sp = slide.shapes.add_shape(shape, left, top, width, height)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def add_logos(slide, sapct=True):
    """Marca SAPCT (sup. direito) e CIMATEC (inf. direito)."""
    if sapct:
        add_pic_fit(slide, LOGO_SAPCT, Inches(11.02), Inches(0.34),
                    Inches(1.78), Inches(1.02), align="right")
    add_pic_fit(slide, LOGO_CIMATEC, Inches(11.55), Inches(6.86),
                Inches(1.28), Inches(0.52), align="right")


def add_footer(slide, number):
    """Rodape: regua + texto do evento + numero do slide."""
    add_rect(slide, Inches(0.62), Inches(7.04), Inches(7.6), Pt(1.6), TEAL)
    add_text(slide, Inches(0.62), Inches(7.10), Inches(7.8), Inches(0.3),
             [[(FOOTER_TXT, 8.5, GRAY)]])
    add_text(slide, Inches(10.6), Inches(7.04), Inches(0.85), Inches(0.32),
             [[("%02d" % number, 10, BLUE, True)]], align=PP_ALIGN.RIGHT)


def add_header(slide, kicker, title):
    """Cabecalho padrao: rotulo (kicker) + titulo + regua de destaque."""
    # rotulo de secao em pilula navy
    pill = add_rect(slide, Inches(0.62), Inches(0.42), Inches(0.22),
                    Inches(0.30), TEAL)
    add_text(slide, Inches(0.95), Inches(0.40), Inches(8.0), Inches(0.34),
             [[(kicker.upper(), 11, TEAL, True)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.60), Inches(0.74), Inches(10.1), Inches(0.92),
             [[(title, 27, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.0)
    add_rect(slide, Inches(0.64), Inches(1.58), Inches(0.95), Pt(3.4), NAVY)


def content_slide(prs, kicker, title, number, sapct=True):
    """Cria um slide de conteudo padrao e devolve o objeto slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # BLANK
    clear_placeholders(slide)
    set_bg(slide, WHITE)
    add_header(slide, kicker, title)
    add_logos(slide, sapct=sapct)
    add_footer(slide, number)
    return slide


def caption(slide, left, top, width, text):
    add_text(slide, left, top, width, Inches(0.3),
             [[(text, 9, GRAY, False, True)]], align=PP_ALIGN.CENTER)


def stat_card(slide, left, top, width, height, big, label, big_color=NAVY,
              fill=LIGHT):
    """Cartao de indicador: numero grande + rotulo."""
    add_rect(slide, left, top, width, height, fill,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    add_rect(slide, left, top, Inches(0.07), height, TEAL,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    add_text(slide, left + Inches(0.18), top + Inches(0.10),
             width - Inches(0.30), Inches(0.50),
             [[(big, 23, big_color, True)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, left + Inches(0.18), top + height - Inches(0.62),
             width - Inches(0.30), Inches(0.55),
             [[(label, 10.5, GRAY)]], anchor=MSO_ANCHOR.TOP, line_spacing=1.02)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ======================================================================
# MONTAGEM DA APRESENTACAO
# ======================================================================
prs = Presentation(TEMPLATE)

# --- mantem o slide 1 (capa do evento) e remove os 2 slides-modelo -----
xml_slides = prs.slides._sldIdLst
slide_ids = list(xml_slides)
for sid in slide_ids[1:]:
    rId = sid.get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(sid)

RID = qn("r:id")

# ----------------------------------------------------------------------
# SLIDE 2 - Titulo do trabalho
# ----------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[0])
clear_placeholders(s)
set_bg(s, WHITE)
# faixa lateral de cor
add_rect(s, 0, 0, Inches(0.34), EMU_H, NAVY)
add_rect(s, Inches(0.34), 0, Inches(0.06), EMU_H, TEAL)
# marca SAPCT grande
add_pic_fit(s, LOGO_SAPCT, Inches(0.95), Inches(0.62), Inches(3.0), Inches(1.55),
            align="left")
add_text(s, Inches(0.98), Inches(2.26), Inches(7.0), Inches(0.32),
         [[("REVISÃO EXPLORATÓRIA DA LITERATURA", 12, TEAL, True)]])
add_text(s, Inches(0.98), Inches(2.66), Inches(11.4), Inches(2.0),
         [[("Aplicação de Quantum Annealing ao Problema do "
            "Caixeiro Viajante", 33, NAVY, True)],
          [("Revisão exploratória da literatura em logística",
            21, BLUE, False, True)]],
         line_spacing=1.05, space_after=6)
add_rect(s, Inches(1.0), Inches(4.30), Inches(1.1), Pt(3.4), TEAL)
# autores
add_text(s, Inches(0.98), Inches(4.52), Inches(11.2), Inches(0.9),
         [[("Renan Cardoso dos Santos", 15, INK, True),
           ("¹    ", 11, TEAL, True),
           ("Valéria Loureiro da Silva", 15, INK, True),
           ("²    ", 11, TEAL, True),
           ("Anderson Rafael Correia B. da Silva", 15, INK, True),
           ("³", 11, TEAL, True)]])
add_text(s, Inches(0.98), Inches(4.92), Inches(11.4), Inches(1.0),
         [[("¹ Mestrando em Gestão de Tecnologia e Inovação "
            "— Universidade SENAI CIMATEC, Salvador-BA", 11, GRAY)],
          [("² ³ Universidade SENAI CIMATEC, Salvador-BA",
            11, GRAY)]], space_after=3)
# palavras-chave
add_rect(s, Inches(0.98), Inches(5.86), Inches(8.7), Inches(0.66), LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
add_text(s, Inches(1.20), Inches(5.86), Inches(8.3), Inches(0.66),
         [[("Palavras-chave:  ", 10.5, TEAL, True),
           ("TSP · Quantum Annealing · Otimização "
            "Combinatória · Logística", 10.5, INK)]],
         anchor=MSO_ANCHOR.MIDDLE)
add_pic_fit(s, LOGO_CIMATEC, Inches(10.0), Inches(5.74), Inches(2.5),
            Inches(0.95), align="right")
add_text(s, Inches(0.98), Inches(6.84), Inches(11.0), Inches(0.4),
         [[("XI Seminário de Avaliação de Pesquisa "
            "Científica e Tecnológica — SAPCT 2026", 10,
            GRAY, False, True)]])
notes(s, "Apresentar o trabalho: revisao exploratoria da literatura sobre o "
         "uso de Quantum Annealing no Problema do Caixeiro Viajante (TSP) "
         "aplicado a logistica. Mestrado Profissional em Gestao de Tecnologia "
         "e Inovacao - SENAI CIMATEC.")

# ----------------------------------------------------------------------
# SLIDE 3 - Roteiro
# ----------------------------------------------------------------------
s = content_slide(prs, "Roteiro", "O que vamos percorrer", 3)
roteiro = [
    ("O desafio", "TSP e a otimização de rotas em logística"),
    ("A promessa", "Computação quântica e Quantum Annealing"),
    ("Objetivo", "A pergunta que orienta a revisão"),
    ("Metodologia", "Como a literatura foi mapeada"),
    ("Resultados", "Panorama, paradigmas e algoritmos"),
    ("Lacunas & rumos", "Oportunidades e próximas etapas"),
]
cw, ch, gx, gy = Inches(3.86), Inches(1.62), Inches(0.30), Inches(0.30)
x0, y0 = Inches(0.66), Inches(2.05)
for i, (t, d) in enumerate(roteiro):
    col, row = i % 3, i // 3
    l = x0 + col * (cw + gx)
    tp = y0 + row * (ch + gy)
    add_rect(s, l, tp, cw, ch, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.07)
    add_rect(s, l, tp, cw, Inches(0.07), TEAL)
    add_text(s, l + Inches(0.26), tp + Inches(0.16), Inches(1.2), Inches(0.5),
             [[("%02d" % (i + 1), 26, RGBColor(0xD3, 0xDD, 0xE7), True)]])
    add_text(s, l + Inches(0.26), tp + Inches(0.62), cw - Inches(0.5),
             Inches(0.4), [[(t, 15, NAVY, True)]])
    add_text(s, l + Inches(0.26), tp + Inches(0.98), cw - Inches(0.5),
             Inches(0.55), [[(d, 11, GRAY)]], line_spacing=1.05)
notes(s, "Roteiro em seis blocos: contextualizar o problema, apresentar a "
         "tecnologia, definir o objetivo, descrever a metodologia, discutir "
         "os resultados e fechar com lacunas e proximos passos.")

# ----------------------------------------------------------------------
# SLIDE 4 - O desafio: TSP & Logistica
# ----------------------------------------------------------------------
s = content_slide(prs, "O Desafio", "TSP: um problema NP-Hard no coração "
                   "da logística", 4)
add_rect(s, Inches(0.66), Inches(1.95), Inches(5.66), Inches(4.05), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_rect(s, Inches(0.66), Inches(1.95), Inches(5.66), Inches(0.07), NAVY)
add_text(s, Inches(0.95), Inches(2.16), Inches(5.1), Inches(0.4),
         [[("O Problema do Caixeiro Viajante", 14, NAVY, True)]])
add_bullets(s, Inches(0.95), Inches(2.62), Inches(5.1), Inches(3.2), [
    ("Determinar a rota de menor custo que visite um conjunto de cidades "
     "exatamente uma vez e retorne à origem.",),
    ("Classificado como NP-Hard — o esforço computacional cresce "
     "explosivamente com o número de cidades.",),
    ("Grande relevância prática: a otimização de rotas de "
     "entrega impacta diretamente custos operacionais e eficiência.",),
], size=13.5, gap=10)

add_rect(s, Inches(6.60), Inches(1.95), Inches(6.07), Inches(4.05), NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_rect(s, Inches(6.60), Inches(1.95), Inches(6.07), Inches(0.07), TEAL)
add_text(s, Inches(6.90), Inches(2.16), Inches(5.5), Inches(0.4),
         [[("Da teoria ao chão logístico: o TSPTW", 14, WHITE, True)]])
add_bullets(s, Inches(6.90), Inches(2.62), Inches(5.5), Inches(3.2), [
    ("O TSP com Janelas de Tempo (TSPTW) incorpora restrições "
     "temporais de atendimento a cada cidade.",),
    ("Aproxima o modelo de cenários reais como distribuição "
     "urbana e roteamento de veículos.",),
    ("Variantes correlatas também relevantes: VRP, CVRP, Job Shop "
     "Scheduling e Bin Packing.",),
], size=13.5, gap=10, color=RGBColor(0xDD, 0xE6, 0xEF), lead_color=STEEL)
add_text(s, Inches(6.90), Inches(5.30), Inches(5.5), Inches(0.6),
         [[("“Métodos clássicos enfrentam limites práticos "
            "à medida que as instâncias crescem.”", 11.5,
            STEEL, False, True)]], line_spacing=1.1)
notes(s, "O TSP e um classico NP-Hard de otimizacao combinatoria. Na "
         "logistica, rotas eficientes reduzem custo e tempo. A variante "
         "TSPTW, com janelas de tempo, aproxima o problema de operacoes "
         "reais de distribuicao urbana e roteamento de veiculos.")

# ----------------------------------------------------------------------
# SLIDE 5 - A promessa: Computacao Quantica & QA
# ----------------------------------------------------------------------
s = content_slide(prs, "A Promessa", "Quantum Annealing: explorar "
                   "soluções de outra forma", 5)
steps = [
    ("Computação quântica",
     "Explora o espaço de soluções de problemas combinatórios "
     "de forma mais eficiente que métodos clássicos."),
    ("Quantum Annealing (QA)",
     "Usa o tunelamento quântico para buscar o estado de menor energia "
     "de um sistema."),
    ("Formulação QUBO",
     "O problema é codificado como Quadratic Unconstrained Binary "
     "Optimization (QUBO / Ising)."),
    ("Hardware D-Wave",
     "QUBO é resolvido em processadores quânticos dedicados, "
     "como os da plataforma D-Wave."),
]
cw = Inches(2.86)
gx = Inches(0.22)
x0 = Inches(0.66)
for i, (t, d) in enumerate(steps):
    l = x0 + i * (cw + gx)
    add_rect(s, l, Inches(2.05), cw, Inches(2.62), CARD,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    circ = add_rect(s, l + Inches(0.24), Inches(2.26), Inches(0.52),
                    Inches(0.52), TEAL, shape=MSO_SHAPE.OVAL)
    add_text(s, l + Inches(0.24), Inches(2.26), Inches(0.52), Inches(0.52),
             [[(str(i + 1), 16, WHITE, True)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, l + Inches(0.24), Inches(2.92), cw - Inches(0.46),
             Inches(0.62), [[(t, 13, NAVY, True)]], line_spacing=1.0)
    add_text(s, l + Inches(0.24), Inches(3.50), cw - Inches(0.46),
             Inches(1.1), [[(d, 10.8, GRAY)]], line_spacing=1.08)
    if i < 3:
        add_text(s, l + cw - Inches(0.02), Inches(3.0), Inches(0.30),
                 Inches(0.4), [[("›", 20, STEEL, True)]],
                 align=PP_ALIGN.CENTER)
# destaque referencia
add_rect(s, Inches(0.66), Inches(4.92), Inches(12.01), Inches(1.08), NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
add_text(s, Inches(1.0), Inches(4.92), Inches(11.4), Inches(1.08),
         [[("Referência-chave   ", 11, STEEL, True),
           ("Salehi, Glos e Miszczak (2022) propuseram três "
            "formulações QUBO para o TSPTW, demonstrando a "
            "viabilidade de codificar variantes realistas do problema "
            "em dispositivos de quantum annealing.", 12.5, WHITE)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
notes(s, "A computacao quantica busca explorar o espaco de solucoes de forma "
         "mais eficiente. O Quantum Annealing usa tunelamento quantico para "
         "encontrar o estado de menor energia; o problema e formulado como "
         "QUBO e resolvido em hardware D-Wave. Salehi et al. (2022) "
         "mostraram a viabilidade para o TSPTW.")

# ----------------------------------------------------------------------
# SLIDE 6 - Objetivo
# ----------------------------------------------------------------------
s = content_slide(prs, "Objetivo", "A pergunta que orienta a revisão", 6)
add_rect(s, Inches(0.66), Inches(2.00), Inches(12.01), Inches(1.30), LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
add_rect(s, Inches(0.66), Inches(2.00), Inches(0.09), Inches(1.30), TEAL)
add_text(s, Inches(1.05), Inches(2.00), Inches(11.2), Inches(1.30),
         [[("Qual o panorama atual da aplicação de algoritmos "
            "quânticos — com ênfase em Quantum Annealing — "
            "ao TSP e suas variantes em logística?", 17, NAVY, True)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
add_text(s, Inches(0.66), Inches(3.62), Inches(8.0), Inches(0.4),
         [[("Para responder, a revisão busca identificar:", 13,
            BLUE, True)]])
goals = [
    ("Volume", "O volume de produção acadêmica sobre o tema."),
    ("Algoritmos", "Os algoritmos e paradigmas quânticos mais utilizados."),
    ("Lacunas", "As lacunas de pesquisa ainda existentes."),
]
cw = Inches(3.86)
gx = Inches(0.215)
for i, (t, d) in enumerate(goals):
    l = Inches(0.66) + i * (cw + gx)
    add_rect(s, l, Inches(4.10), cw, Inches(1.86), CARD,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
    add_rect(s, l, Inches(4.10), cw, Inches(0.07), TEAL)
    add_text(s, l + Inches(0.28), Inches(4.30), cw - Inches(0.5), Inches(0.45),
             [[(t.upper(), 12.5, TEAL, True)]])
    add_text(s, l + Inches(0.28), Inches(4.74), cw - Inches(0.5), Inches(1.1),
             [[(d, 12.5, INK)]], line_spacing=1.12)
notes(s, "Objetivo: realizar uma revisao exploratoria da literatura sobre a "
         "aplicacao de algoritmos quanticos, com enfase em Quantum Annealing, "
         "ao TSP e variantes em logistica - identificando volume de producao, "
         "algoritmos mais usados e lacunas de pesquisa.")

# ----------------------------------------------------------------------
# SLIDE 7 - Metodologia
# ----------------------------------------------------------------------
s = content_slide(prs, "Metodologia", "Como a literatura foi mapeada", 7)
add_text(s, Inches(0.66), Inches(1.92), Inches(6.0), Inches(0.4),
         [[("Pesquisa bibliográfica exploratória em três eixos "
            "temáticos", 13, BLUE, True)]])
eixos = [
    ("1 · O problema", "TSP e variantes — VRP, CVRP, Job Shop "
     "Scheduling e Bin Packing."),
    ("2 · A tecnologia", "Paradigmas quânticos — QA, QAOA, VQE, "
     "Grover e formulações QUBO/Ising."),
    ("3 · A aplicação", "Logística, gestão da cadeia de "
     "suprimentos e otimização de rotas."),
]
for i, (t, d) in enumerate(eixos):
    tp = Inches(2.32) + i * Inches(0.83)
    add_rect(s, Inches(0.66), tp, Inches(5.95), Inches(0.72), CARD,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    add_rect(s, Inches(0.66), tp, Inches(0.07), Inches(0.72), TEAL)
    add_text(s, Inches(0.92), tp, Inches(1.78), Inches(0.72),
             [[(t, 12, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.0)
    add_text(s, Inches(2.78), tp, Inches(3.70), Inches(0.72),
             [[(d, 10.6, GRAY)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
add_bullets(s, Inches(0.66), Inches(4.95), Inches(6.0), Inches(2.0), [
    ("26 strings de busca (operador AND) na base Lens.org — filtro "
     "Scholarly Works.",),
    ("Prioridades: 8 Alta · 10 Média · 8 Baixa; em duas fases "
     "(16 strings + 10 derivadas de Phillipson, 2025).",),
    ("Deduplicação em duas etapas: por DOI e, depois, por "
     "título normalizado.",),
], size=11.5, gap=7)
# grafico de strings
add_rect(s, Inches(6.78), Inches(1.92), Inches(5.92), Inches(4.42), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_pic_fit(s, os.path.join(IMG, "06_strings_busca_01_volume_strings.png"),
            Inches(6.86), Inches(2.00), Inches(5.76), Inches(3.95))
caption(s, Inches(6.78), Inches(5.96), Inches(5.92),
        "Volume de artigos únicos recuperados por string de busca.")
notes(s, "A metodologia cruza tres eixos - problema, tecnologia e aplicacao - "
         "em 26 strings de busca executadas no Lens.org. As strings foram "
         "priorizadas e aplicadas em duas fases. A deduplicacao por DOI e "
         "titulo garante a consistencia do corpus.")

# ----------------------------------------------------------------------
# SLIDE 8 - Resultados: panorama / expansao
# ----------------------------------------------------------------------
s = content_slide(prs, "Resultados · 1", "Um campo de pesquisa em "
                   "franca expansão", 8)
cards = [
    ("5.320", "Registros brutos recuperados nas 26 strings"),
    ("3.696", "Artigos únicos após deduplicação"),
    ("30,5%", "Duplicatas removidas (1.624 registros)"),
    ("61,7%", "Publicados a partir de 2020 — pico em 2025"),
]
cw = Inches(2.86)
gx = Inches(0.22)
for i, (b, lab) in enumerate(cards):
    l = Inches(0.66) + i * (cw + gx)
    stat_card(s, l, Inches(1.96), cw, Inches(1.30), b, lab)
add_rect(s, Inches(0.66), Inches(3.46), Inches(12.01), Inches(2.92), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_pic_fit(s, os.path.join(IMG, "01_visao_geral_01_publicacoes_tempo.png"),
            Inches(0.80), Inches(3.54), Inches(11.73), Inches(2.62))
caption(s, Inches(0.66), Inches(6.30), Inches(12.01),
        "Publicações ao longo do tempo, por tipo de documento "
        "— forte aceleração a partir de 2020.")
notes(s, "Dos 5.320 registros brutos, a deduplicacao resultou em 3.696 "
         "artigos unicos (30,5% de duplicatas). 61,7% foram publicados a "
         "partir de 2020, com pico em 2025: um campo em franca expansao.")

# ----------------------------------------------------------------------
# SLIDE 9 - Resultados: alcance global & ciencia aberta
# ----------------------------------------------------------------------
s = content_slide(prs, "Resultados · 2", "Alcance global e ciência "
                   "aberta", 9)
for i, (b, lab) in enumerate([
        ("49", "Países representados na produção"),
        ("1.409", "Periódicos e conferências distintos"),
        ("82,8%", "Artigos em acesso aberto (Open Access)")]):
    l = Inches(0.66) + i * Inches(4.07)
    stat_card(s, l, Inches(1.96), Inches(3.86), Inches(1.20), b, lab)
add_rect(s, Inches(0.66), Inches(3.34), Inches(6.94), Inches(3.04), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_pic_fit(s, os.path.join(IMG, "05_geografia_01_mapa_mundial.png"),
            Inches(0.74), Inches(3.42), Inches(6.78), Inches(2.62))
caption(s, Inches(0.66), Inches(6.04), Inches(6.94),
        "Distribuição geográfica das publicações.")
add_rect(s, Inches(7.74), Inches(3.34), Inches(4.93), Inches(3.04), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_pic_fit(s, os.path.join(IMG, "01_visao_geral_03_open_access.png"),
            Inches(7.82), Inches(3.42), Inches(4.77), Inches(2.62))
caption(s, Inches(7.74), Inches(6.04), Inches(4.93),
        "Open Access vs. acesso restrito.")
notes(s, "A producao e globalmente distribuida: 49 paises e 1.409 periodicos "
         "e conferencias distintos. 82,8% dos artigos estao em acesso aberto, "
         "o que favorece a reprodutibilidade e a replicabilidade da pesquisa.")

# ----------------------------------------------------------------------
# SLIDE 10 - Resultados: paradigma dominante (QA)
# ----------------------------------------------------------------------
s = content_slide(prs, "Resultados · 3", "O paradigma dominante: "
                   "Quantum Annealing", 10)
add_bullets(s, Inches(0.66), Inches(2.02), Inches(5.9), Inches(2.4), [
    ("O Quantum Annealing é o paradigma dominante: mais de 1.200 "
     "artigos nas strings diretamente relacionadas.",),
    ("Hardware de referência — D-Wave Advantage: 5.640 qubits e "
     "conectividade Pegasus.",),
    ("Algoritmos variacionais gate-based (QAOA, VQE) somam ≈ 150 "
     "artigos aplicados ao TSP — uma lacuna evidente.",),
], size=12.5, gap=9)
stat_card(s, Inches(0.66), Inches(4.60), Inches(2.86), Inches(1.74),
          "> 1.200", "Artigos em Quantum Annealing", big_color=NAVY)
stat_card(s, Inches(3.74), Inches(4.60), Inches(2.82), Inches(1.74),
          "≈ 150", "Artigos QAOA / VQE aplicados ao TSP",
          big_color=BLUE)
add_rect(s, Inches(6.78), Inches(2.02), Inches(5.89), Inches(2.20), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
add_pic_fit(s, os.path.join(IMG, "07_algoritmos_01_paradigma.png"),
            Inches(6.86), Inches(2.08), Inches(5.73), Inches(1.92))
add_rect(s, Inches(6.78), Inches(4.34), Inches(5.89), Inches(2.00), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
add_pic_fit(s, os.path.join(IMG, "07_algoritmos_04_algoritmos_frequentes.png"),
            Inches(6.86), Inches(4.40), Inches(5.73), Inches(1.74))
caption(s, Inches(6.78), Inches(6.28), Inches(5.89),
        "Amostra de artigos classificados — paradigma (acima) e "
        "algoritmos mais frequentes (abaixo).")
notes(s, "O Quantum Annealing confirma-se como paradigma dominante, com mais "
         "de 1.200 artigos. O hardware de referencia e o D-Wave Advantage "
         "(5.640 qubits, Pegasus). Algoritmos gate-based como QAOA e VQE "
         "somam apenas cerca de 150 artigos aplicados ao TSP - uma lacuna.")

# ----------------------------------------------------------------------
# SLIDE 11 - Resultados: taxonomia & abordagens hibridas
# ----------------------------------------------------------------------
s = content_slide(prs, "Resultados · 4", "Taxonomia de algoritmos e "
                   "abordagens híbridas", 11)
add_text(s, Inches(0.66), Inches(1.94), Inches(6.2), Inches(0.4),
         [[("Mais de 40 algoritmos catalogados → 6 categorias "
            "principais", 13, BLUE, True)]])
cats = [
    ("Quantum Annealing", "D-Wave QA, Hybrid Solver, Reverse Annealing"),
    ("Variacionais", "QAOA, VQE, F-VQE"),
    ("Exatos", "Grover, QPE"),
    ("Híbridos", "Decomposição + QA, path-slicing"),
    ("QML", "Quantum RL, QACO"),
    ("Quantum-inspired", "Heurísticas inspiradas no quântico"),
]
cw, ch, gx, gy = Inches(3.00), Inches(0.84), Inches(0.16), Inches(0.14)
for i, (t, d) in enumerate(cats):
    col, row = i % 2, i // 2
    l = Inches(0.66) + col * (cw + gx)
    tp = Inches(2.36) + row * (ch + gy)
    add_rect(s, l, tp, cw, ch, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.10)
    add_rect(s, l, tp, Inches(0.07), ch, TEAL)
    add_text(s, l + Inches(0.20), tp + Inches(0.09), cw - Inches(0.3),
             Inches(0.34), [[(t, 11.5, NAVY, True)]])
    add_text(s, l + Inches(0.20), tp + Inches(0.41), cw - Inches(0.3),
             Inches(0.4), [[(d, 9, GRAY)]], line_spacing=1.0)
add_rect(s, Inches(0.66), Inches(5.42), Inches(6.16), Inches(0.92), NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
add_text(s, Inches(0.92), Inches(5.42), Inches(5.7), Inches(0.92),
         [[("A maioria das implementações adota abordagens "
            "híbridas", 12, WHITE, True)],
          [("processamento quântico combinado com solvers "
            "clássicos.", 10.5, STEEL)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05, space_after=2)
add_rect(s, Inches(7.00), Inches(1.94), Inches(5.67), Inches(4.40), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_pic_fit(s, os.path.join(IMG, "07_algoritmos_02_abordagem.png"),
            Inches(7.08), Inches(2.02), Inches(5.51), Inches(3.90))
caption(s, Inches(7.00), Inches(5.94), Inches(5.67),
        "Distribuição por abordagem na amostra classificada — "
        "predomínio de soluções híbridas.")
notes(s, "Foram catalogados mais de 40 algoritmos quanticos, agrupados em "
         "seis categorias: Quantum Annealing, variacionais, exatos, hibridos, "
         "QML e quantum-inspired. A maioria das implementacoes combina "
         "processamento quantico com solvers classicos - abordagens hibridas.")

# ----------------------------------------------------------------------
# SLIDE 12 - Lacunas de pesquisa
# ----------------------------------------------------------------------
s = content_slide(prs, "Lacunas", "Onde estão as oportunidades de "
                   "pesquisa", 12)
add_text(s, Inches(0.66), Inches(1.92), Inches(11.8), Inches(0.4),
         [[("Quatro lacunas se destacam — e apontam direções "
            "promissoras para trabalhos futuros:", 13, BLUE, True)]])
lac = [
    ("Gate-based", "Escassas",
     "Falta de soluções QAOA/VQE com escalabilidade demonstrada "
     "para o TSP."),
    ("QRL p/ roteamento", "21 artigos",
     "Quantum Reinforcement Learning para roteamento é área "
     "emergente e pouco explorada."),
    ("CVRP com QA", "18 artigos",
     "Aplicação logística direta do Quantum Annealing ao CVRP "
     "ainda pouco investigada."),
    ("Dados reais", "Limitada",
     "Validação com dados logísticos reais — predominam "
     "benchmarks acadêmicos."),
]
cw = Inches(2.86)
gx = Inches(0.22)
for i, (t, big, d) in enumerate(lac):
    l = Inches(0.66) + i * (cw + gx)
    add_rect(s, l, Inches(2.46), cw, Inches(3.30), CARD,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    add_rect(s, l, Inches(2.46), cw, Inches(0.86), NAVY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.22)
    add_rect(s, l, Inches(2.92), cw, Inches(0.40), NAVY)
    add_text(s, l, Inches(2.46), cw, Inches(0.86),
             [[(t, 13, WHITE, True)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, l + Inches(0.20), Inches(3.46), cw - Inches(0.4), Inches(0.6),
             [[(big, 21, TEAL, True)]], align=PP_ALIGN.CENTER)
    add_text(s, l + Inches(0.26), Inches(4.12), cw - Inches(0.52),
             Inches(1.5), [[(d, 11.5, INK)]], line_spacing=1.16,
             align=PP_ALIGN.CENTER)
add_text(s, Inches(0.66), Inches(5.96), Inches(12.0), Inches(0.4),
         [[("Também persiste o desafio de demonstrar escalabilidade "
            "frente às limitações do hardware NISQ.", 11,
            GRAY, False, True)]], align=PP_ALIGN.CENTER)
notes(s, "Quatro lacunas principais: ausencia de solucoes gate-based "
         "escalaveis para o TSP; apenas 21 artigos de Quantum Reinforcement "
         "Learning para roteamento; somente 18 artigos de CVRP com QA; e "
         "validacao limitada com dados logisticos reais. Sao oportunidades "
         "concretas de contribuicao cientifica.")

# ----------------------------------------------------------------------
# SLIDE 13 - Consideracoes finais & proximas etapas
# ----------------------------------------------------------------------
s = content_slide(prs, "Conclusão", "Considerações finais e "
                   "próximas etapas", 13)
add_rect(s, Inches(0.66), Inches(1.98), Inches(5.96), Inches(4.42), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_rect(s, Inches(0.66), Inches(1.98), Inches(5.96), Inches(0.07), NAVY)
add_text(s, Inches(0.95), Inches(2.18), Inches(5.4), Inches(0.4),
         [[("O que a revisão mostrou", 14, NAVY, True)]])
add_bullets(s, Inches(0.95), Inches(2.66), Inches(5.4), Inches(3.6), [
    ("Campo em rápida expansão, sustentado por uma base robusta "
     "de 3.696 artigos únicos.",),
    ("O Quantum Annealing é o paradigma mais maduro e acessível "
     "para otimização combinatória.",),
    ("Algoritmos variacionais e de aprendizado de máquina quântico "
     "são fronteiras ainda pouco exploradas.",),
    ("As lacunas identificadas são oportunidades concretas de "
     "contribuição científica.",),
], size=12.5, gap=11)
add_rect(s, Inches(6.72), Inches(1.98), Inches(5.95), Inches(4.42), NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_rect(s, Inches(6.72), Inches(1.98), Inches(5.95), Inches(0.07), TEAL)
add_text(s, Inches(7.02), Inches(2.18), Inches(5.4), Inches(0.4),
         [[("Próximas etapas da pesquisa", 14, WHITE, True)]])
steps = [
    "Fundamentação teórica aprofundada do tema.",
    "Seleção e reprodução de algoritmos em simuladores "
    "quânticos.",
    "Comparação experimental entre abordagens clássicas e "
    "quânticas para instâncias do TSP em contexto logístico.",
]
for i, t in enumerate(steps):
    tp = Inches(2.74) + i * Inches(1.16)
    add_rect(s, Inches(7.02), tp, Inches(0.54), Inches(0.54), TEAL,
             shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(7.02), tp, Inches(0.54), Inches(0.54),
             [[(str(i + 1), 16, WHITE, True)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.74), tp - Inches(0.04), Inches(4.7), Inches(1.1),
             [[(t, 12.5, RGBColor(0xE2, 0xE9, 0xF1))]],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.14)
notes(s, "A revisao mostra um campo em rapida expansao. O QA e o paradigma "
         "mais maduro; variacionais e QML sao fronteiras abertas. As "
         "proximas etapas: fundamentacao teorica, reproducao de algoritmos "
         "em simuladores e comparacao experimental classico vs. quantico "
         "para o TSP logistico.")

# ----------------------------------------------------------------------
# SLIDE 14 - Referencias & Agradecimentos
# ----------------------------------------------------------------------
s = content_slide(prs, "Referências", "Referências e agradecimentos",
                   14)
refs = [
    "APPLEGATE, D. L. et al. The Traveling Salesman Problem: a computational "
    "study. Princeton: Princeton University Press, 2006.",
    "KADOWAKI, T.; NISHIMORI, H. Quantum annealing in the Transverse Ising "
    "model. Physical Review E, v. 58, n. 5, p. 5355-5363, 1998.",
    "SALEHI, Ö.; GLOS, A.; MISZCZAK, J. A. Unconstrained binary models "
    "of the travelling salesman problem variants for quantum optimization. "
    "Quantum Information Processing, v. 21, n. 67, 2022.",
    "PHILLIPSON, F. Quantum Computing in Logistics and Supply Chain "
    "Management: an overview. Maastricht University / TNO, 2025.",
    "FARHI, E.; GOLDSTONE, J.; GUTMANN, S. A quantum approximate "
    "optimization algorithm. arXiv preprint arXiv:1411.4028, 2014.",
]
add_rect(s, Inches(0.66), Inches(1.96), Inches(7.55), Inches(4.42), CARD,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
add_rect(s, Inches(0.66), Inches(1.96), Inches(7.55), Inches(0.07), NAVY)
add_text(s, Inches(0.94), Inches(2.14), Inches(7.0), Inches(0.4),
         [[("Principais referências", 13, NAVY, True)]])
tb = s.shapes.add_textbox(Inches(0.94), Inches(2.58), Inches(7.05),
                          Inches(3.7))
tf = tb.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
for i, r in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.line_spacing = 1.1
    p.space_after = Pt(9)
    num = p.add_run()
    num.text = "%d.  " % (i + 1)
    _style_run(num, 10.5, TEAL, True)
    body = p.add_run()
    body.text = r
    _style_run(body, 10.5, INK)
# agradecimentos
add_rect(s, Inches(8.40), Inches(1.96), Inches(4.27), Inches(4.42), NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
add_rect(s, Inches(8.40), Inches(1.96), Inches(4.27), Inches(0.07), TEAL)
add_text(s, Inches(8.68), Inches(2.14), Inches(3.8), Inches(0.4),
         [[("Agradecimentos", 13, WHITE, True)]])
add_text(s, Inches(8.68), Inches(2.62), Inches(3.78), Inches(3.6),
         [[("Trabalho parcialmente financiado pelo projeto QuIIN FCRH "
            "Mestrado TQ, suportado pelo QuIIN — Quantum Industrial "
            "Innovation, Centro de Competência EMBRAPII CIMATEC em "
            "Tecnologias Quânticas.", 10.5,
            RGBColor(0xDD, 0xE6, 0xEF))],
          [("Recursos do PPI IoT/Manufatura 4.0 do MCTI, via Termo de "
            "Cooperação 053/2023 com a EMBRAPII.", 10.5,
            RGBColor(0xDD, 0xE6, 0xEF))],
          [("Apoio adicional do CNPq — Conselho Nacional de "
            "Desenvolvimento Científico e Tecnológico.", 10.5,
            RGBColor(0xDD, 0xE6, 0xEF))]],
         line_spacing=1.16, space_after=8)
notes(s, "Referencias centrais do trabalho e agradecimentos as agencias e "
         "projetos de fomento: QuIIN, EMBRAPII CIMATEC, MCTI e CNPq.")

# ----------------------------------------------------------------------
# SLIDE 15 - Encerramento
# ----------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[0])
clear_placeholders(s)
set_bg(s, NAVY)
add_rect(s, 0, Inches(6.86), EMU_W, Inches(0.64), TEAL)
add_pic_fit(s, LOGO_SAPCT, Inches(0.95), Inches(0.78), Inches(2.7),
            Inches(1.45), align="left")
add_text(s, Inches(0.98), Inches(2.70), Inches(10.0), Inches(1.3),
         [[("Obrigado!", 50, WHITE, True)]])
add_rect(s, Inches(1.02), Inches(3.86), Inches(1.15), Pt(4), TEAL)
add_text(s, Inches(0.98), Inches(4.06), Inches(11.0), Inches(0.9),
         [[("Aplicação de Quantum Annealing ao Problema do "
            "Caixeiro Viajante:", 16, STEEL, True)],
          [("revisão exploratória da literatura em logística",
            14, RGBColor(0xB9, 0xC6, 0xD6), False, True)]],
         line_spacing=1.1, space_after=4)
add_text(s, Inches(0.98), Inches(5.18), Inches(11.0), Inches(0.5),
         [[("Renan Cardoso dos Santos", 14, WHITE, True),
           ("   ·   Valéria Loureiro da Silva   ·   "
            "Anderson Rafael Correia B. da Silva", 12.5,
            RGBColor(0xB9, 0xC6, 0xD6))]])
add_text(s, Inches(0.98), Inches(5.66), Inches(11.0), Inches(0.5),
         [[("✉  ", 12, TEAL, True),
           ("renan.santos@fbter.org.br", 12.5, STEEL),
           ("        Universidade SENAI CIMATEC · Salvador-BA",
            11.5, RGBColor(0x8F, 0x9F, 0xB3))]])
add_text(s, Inches(0.98), Inches(6.98), Inches(9.0), Inches(0.4),
         [[("XI SAPCT — SENAI CIMATEC · 28 e 29 de Maio de 2026",
            10.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
# marca CIMATEC em cartao branco (logo possui fundo branco)
add_rect(s, Inches(10.62), Inches(6.92), Inches(2.10), Inches(0.50), WHITE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
add_pic_fit(s, LOGO_CIMATEC, Inches(10.74), Inches(6.96), Inches(1.86),
            Inches(0.42), align="center")
notes(s, "Encerramento e abertura para perguntas. Contato do autor para "
         "interessados na continuidade da pesquisa.")

# ----------------------------------------------------------------------
prs.save(OUTPUT)
print("OK ->", OUTPUT)
print("Total de slides:", len(prs.slides._sldIdLst))
